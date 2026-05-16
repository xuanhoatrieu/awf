"""
PPTX Generator - Core engine for awf-pptx skill
Supports bg1 (background images) and bg2 (programmatic MiniMax style)
8 slide types: Title, Objectives, Agenda, Divider, Content, Summary, Exercise, Closing
"""
import os, sys, re, json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
from typing import List, Dict, Any, Optional

if sys.platform == "win32":
    try: sys.stdout.reconfigure(encoding="utf-8")
    except: pass

# ============================================================
# THEMES
# ============================================================
THEME_BG1 = {
    "title_font": "Georgia", "body_font": "Arial",
    "title_color": RGBColor(255,255,255),
    "body_color": RGBColor(58,102,77),
    "heading_color": RGBColor(58,102,77),
    "badge_fill": "FFFFFF", "badge_text_color": "3A664D",
}
THEME_BG2 = {
    "title_font": "Georgia", "body_font": "Calibri",
    "primary": "264653", "secondary": "2a9d8f", "accent": "e9c46a",
    "light": "f4a261", "bg": "f2e9e4", "alert": "e76f51",
    "badge_fill": "2a9d8f", "badge_text_color": "FFFFFF",
}

def hex2rgb(h):
    return RGBColor(int(h[:2],16), int(h[2:4],16), int(h[4:6],16))

# ============================================================
# PPTX GENERATOR CLASS
# ============================================================
class PptxGenerator:
    def __init__(self, mode="bg1", bg_dir=None):
        self.mode = mode
        self.bg_dir = bg_dir  # pptx/temp_background/
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)
        self.slide_count = 0

    def _blank_layout(self):
        for layout in self.prs.slide_layouts:
            if len(list(layout.placeholders)) == 0:
                return layout
        return self.prs.slide_layouts[-1]

    def _add_slide(self):
        slide = self.prs.slides.add_slide(self._blank_layout())
        # Clear placeholders
        for shape in list(slide.shapes):
            try:
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    shape._element.getparent().remove(shape._element)
            except: pass
        self.slide_count += 1
        return slide

    def _add_bg(self, slide, is_title=False):
        if self.mode == "bg1" and self.bg_dir:
            fn = "1.png" if is_title else "2.png"
            p = os.path.join(self.bg_dir, fn)
            if os.path.exists(p):
                pic = slide.shapes.add_picture(p, Inches(0), Inches(0), Inches(10), Inches(5.625))
                slide.shapes._spTree.insert(2, pic._element)
        elif self.mode == "bg2":
            t = THEME_BG2
            color = t["primary"] if is_title else t["bg"]
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = hex2rgb(color)

    def _page_badge(self, slide, num):
        t = THEME_BG2 if self.mode == "bg2" else THEME_BG1
        x, y, s = Inches(9.3), Inches(5.1), Inches(0.4)
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, s, s)
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex2rgb(t["badge_fill"])
        if self.mode == "bg1":
            shape.fill.fore_color.brightness = 0.5  # semi-transparent effect
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(num)
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = hex2rgb(t["badge_text_color"])
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)

    def _notes(self, slide, text):
        if text:
            slide.notes_slide.notes_text_frame.text = text

    # ---- SLIDE TYPE BUILDERS ----

    def add_title_slide(self, title, subtitle="", note=""):
        slide = self._add_slide()
        self._add_bg(slide, is_title=True)
        t = THEME_BG2 if self.mode == "bg2" else THEME_BG1
        tf_font = t.get("title_font", "Georgia")
        # Title
        sh = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(1.5))
        p = sh.text_frame.paragraphs[0]
        p.text = title
        p.font.name = tf_font; p.font.size = Pt(44); p.font.bold = True
        p.font.color.rgb = RGBColor(255,255,255); p.alignment = PP_ALIGN.CENTER
        # Subtitle
        if subtitle:
            sh2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
            p2 = sh2.text_frame.paragraphs[0]
            p2.text = subtitle; p2.font.name = t.get("body_font","Arial")
            p2.font.size = Pt(22); p2.font.color.rgb = RGBColor(255,255,255)
            p2.alignment = PP_ALIGN.CENTER
        self._notes(slide, note)
        return slide

    def add_objectives_slide(self, title, items, note=""):
        slide = self._add_slide()
        self._add_bg(slide, False)
        self._slide_title(slide, title if title else "Muc tieu bai hoc")
        self._bullet_list(slide, items, start_y=1.2)
        self._page_badge(slide, self.slide_count)
        self._notes(slide, note)
        return slide

    def add_agenda_slide(self, title, sections, note=""):
        slide = self._add_slide()
        self._add_bg(slide, False)
        self._slide_title(slide, title if title else "Noi dung bai hoc")
        items = [f"{i+1}. {s}" for i, s in enumerate(sections)]
        self._bullet_list(slide, items, start_y=1.2)
        self._page_badge(slide, self.slide_count)
        self._notes(slide, note)
        return slide

    def add_divider_slide(self, section_num, section_title, note=""):
        slide = self._add_slide()
        if self.mode == "bg1" and self.bg_dir:
            self._add_bg(slide, False)
        elif self.mode == "bg2":
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = hex2rgb(THEME_BG2["secondary"])
        # Big number
        sh = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.5))
        p = sh.text_frame.paragraphs[0]
        p.text = f"{section_num:02d}"
        p.font.size = Pt(72); p.font.bold = True
        p.font.color.rgb = hex2rgb(THEME_BG2["accent"]) if self.mode == "bg2" else RGBColor(255,255,255)
        p.alignment = PP_ALIGN.CENTER
        # Title
        sh2 = slide.shapes.add_textbox(Inches(1), Inches(3.0), Inches(8), Inches(1))
        p2 = sh2.text_frame.paragraphs[0]
        p2.text = section_title
        p2.font.name = "Georgia"; p2.font.size = Pt(36); p2.font.bold = True
        p2.font.color.rgb = RGBColor(255,255,255); p2.alignment = PP_ALIGN.CENTER
        self._notes(slide, note)
        return slide

    def add_content_slide(self, title, bullets, image_path=None, note="", subtype="text"):
        slide = self._add_slide()
        self._add_bg(slide, False)
        self._slide_title(slide, title)
        has_img = image_path and os.path.exists(image_path)
        w = Inches(4.5) if has_img else Inches(9)
        self._bullet_list(slide, bullets, start_y=1.2, width=w)
        if has_img:
            try:
                slide.shapes.add_picture(image_path, Inches(5.2), Inches(1.0), height=Inches(4.5))
            except Exception as e:
                print(f"[WARN] Image: {e}")
        self._page_badge(slide, self.slide_count)
        self._notes(slide, note)
        return slide

    def add_summary_slide(self, title, points, note=""):
        slide = self._add_slide()
        self._add_bg(slide, False)
        self._slide_title(slide, title if title else "Tong ket")
        self._bullet_list(slide, points, start_y=1.2)
        self._page_badge(slide, self.slide_count)
        self._notes(slide, note)
        return slide

    def add_exercise_slide(self, title, exercises, note=""):
        slide = self._add_slide()
        self._add_bg(slide, False)
        self._slide_title(slide, title if title else "Thao luan & Bai tap")
        # Handle dicts or strings in exercises
        items = []
        for i, ex in enumerate(exercises):
            if isinstance(ex, dict):
                items.append({"point": f"Exercise {i+1}: {ex.get('point','')}", "description": ex.get("description","")})
            else:
                items.append(f"Exercise {i+1}: {ex}")
        self._bullet_list(slide, items, start_y=1.2)
        self._page_badge(slide, self.slide_count)
        self._notes(slide, note)
        return slide

    def add_closing_slide(self, title, next_lesson="", note=""):
        slide = self._add_slide()
        # bg1: closing dung anh 2.png (khong phai 1.png)
        if self.mode == "bg1":
            self._add_bg(slide, is_title=False)  # 2.png
        else:
            self._add_bg(slide, is_title=True)   # bg2: primary color
        sh = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(1.5))
        p = sh.text_frame.paragraphs[0]
        p.text = title if title else "Cam on!"
        p.font.name = "Georgia"; p.font.size = Pt(52)
        p.font.bold = True; p.font.color.rgb = RGBColor(255,255,255); p.alignment = PP_ALIGN.CENTER
        if next_lesson:
            sh2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
            p2 = sh2.text_frame.paragraphs[0]
            prefix = "Next Lesson:" if "Lesson" in next_lesson or "!" in title else "Bai tiep theo:"
            p2.text = f"{prefix} {next_lesson}"
            p2.font.size = Pt(20); p2.font.color.rgb = RGBColor(255,255,255)
            p2.alignment = PP_ALIGN.CENTER
        self._notes(slide, note)
        return slide

    # ---- HELPERS ----

    def _slide_title(self, slide, text):
        t = THEME_BG2 if self.mode == "bg2" else THEME_BG1
        sh = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        p = sh.text_frame.paragraphs[0]
        p.text = text; p.font.name = t.get("title_font", "Georgia")
        p.font.size = Pt(28); p.font.bold = True; p.alignment = PP_ALIGN.CENTER
        if self.mode == "bg1":
            p.font.color.rgb = RGBColor(255,255,255)
        else:
            p.font.color.rgb = hex2rgb(t["primary"])

    def _bullet_list(self, slide, items, start_y=1.2, width=None):
        if not items: return
        t = THEME_BG2 if self.mode == "bg2" else THEME_BG1
        w = width or Inches(9)
        sh = slide.shapes.add_textbox(Inches(0.5), Inches(start_y), w, Inches(4.0))
        tf = sh.text_frame; tf.clear(); tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if isinstance(item, dict):
                point = item.get("point", "")
                desc = item.get("description", "")
                p.text = point; p.font.bold = True; p.font.size = Pt(22)
                if desc:
                    pd = tf.add_paragraph()
                    pd.text = desc; pd.font.bold = False; pd.font.size = Pt(18)
                    pd.level = 1; pd.space_after = Pt(8)
                    pd.font.color.rgb = t.get("body_color") or hex2rgb(t.get("primary","264653"))
                    pd.font.name = t.get("body_font","Arial")
            else:
                p.text = str(item); p.font.bold = False; p.font.size = Pt(20)
            p.font.name = t.get("body_font","Arial"); p.space_after = Pt(6)
            p.font.color.rgb = t.get("body_color") or hex2rgb(t.get("primary","264653"))

    # ---- AUDIO EMBEDDING ----

    def add_audio(self, slide, audio_path):
        """Add audio with auto-play. Reuses proven pattern from pptx_service.py."""
        try:
            from pptx.opc.package import Part
            from pptx.opc.packuri import PackURI
            ext = os.path.splitext(audio_path)[1].lower()
            mime = {'.wav':'audio/wav','.mp3':'audio/mpeg','.m4a':'audio/mp4'}.get(ext)
            if not mime: return
            with open(audio_path,'rb') as f: data = f.read()
            sp = slide.part
            m = re.search(r'slide(\d+)', str(sp.partname))
            si = int(m.group(1))-1 if m else 0
            existing = {str(p.partname) for p in sp.package.iter_parts()}
            mi = 1
            while f"/ppt/media/audio_s{si}_{mi}{ext}" in existing: mi += 1
            pn = f"/ppt/media/audio_s{si}_{mi}{ext}"
            ap = Part(PackURI(pn), mime, sp.package, data)
            rA = sp.relate_to(ap,'http://schemas.openxmlformats.org/officeDocument/2006/relationships/audio')
            rM = sp.relate_to(ap,'http://schemas.microsoft.com/office/2007/relationships/media')
            sid = max((s.shape_id for s in slide.shapes), default=0) + 1
            L,T,W,H = int(-0.6*914400),0,int(0.4*914400),int(0.4*914400)
            # Try speaker icon
            icon_dir = os.path.join(os.path.dirname(__file__), '..', '..', '..', '..', 'pptx', 'bg1')
            icon_path = os.path.join(icon_dir, 'speaker_icon.png')
            rI = ""
            if os.path.exists(icon_path):
                with open(icon_path,'rb') as f: idata = f.read()
                ip = Part(PackURI(f"/ppt/media/spk_s{si}_{mi}.png"),'image/png',sp.package,idata)
                rI = sp.relate_to(ip,'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image')
            xml = f'''<p:pic xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
               xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main">
            <p:nvPicPr><p:cNvPr id="{sid}" name="Audio {sid}">
              <a:hlinkClick r:id="" action="ppaction://media"/></p:cNvPr>
            <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
            <p:nvPr><a:audioFile r:link="{rA}"/>
              <p:extLst><p:ext uri="{{DAA4B4D4-6D71-4841-9C94-3DE7FCFB9230}}">
                <p14:media r:embed="{rM}"/></p:ext></p:extLst></p:nvPr></p:nvPicPr>
            <p:blipFill><a:blip r:embed="{rI}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>
            <p:spPr><a:xfrm><a:off x="{L}" y="{T}"/><a:ext cx="{W}" cy="{H}"/></a:xfrm>
              <a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr></p:pic>'''
            elem = etree.fromstring(xml.strip().encode('utf-8'))
            slide._element.find(qn('p:cSld')).find(qn('p:spTree')).append(elem)
            # Auto-play timing
            txml = f'''<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
            <p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
            <p:childTnLst><p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst><p:par>
            <p:cTn id="3" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>
            <p:childTnLst><p:par><p:cTn id="4" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>
            <p:childTnLst><p:par><p:cTn id="5" presetID="1" presetClass="mediacall" presetSubtype="0" fill="hold" nodeType="afterEffect">
            <p:stCondLst><p:cond delay="0"/></p:stCondLst>
            <p:childTnLst><p:cmd type="call" cmd="playFrom(0.0)"><p:cBhvr>
            <p:cTn id="6" dur="1" fill="hold"/><p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>
            </p:cBhvr></p:cmd></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par>
            </p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:seq>
            </p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>'''
            te = parse_xml(txml)
            old = slide._element.find(qn('p:timing'))
            if old is not None: slide._element.remove(old)
            slide._element.append(te)
            print(f"[PPTX] Audio injected: slide {si}")
        except Exception as e:
            print(f"[WARN] Audio embed failed: {e}")

    def save(self, path):
        os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
        self.prs.save(path)
        print(f"[PPTX] Saved: {path} ({self.slide_count} slides)")


# ============================================================
# BUILD FROM JSON SCRIPT
# ============================================================
def build_from_script(script_path, output_path, mode="bg1", bg_dir=None, audio_dir=None):
    """
    Build PPTX from a JSON script file.

    Script format: list of slide dicts, each with:
      - type: title/objectives/agenda/divider/content/summary/exercise/closing
      - title, subtitle, items, bullets, section_num, section_title, next_lesson
      - note: speaker note text
      - image: path to image file (for content slides)
      - audio: path to audio file (auto-play)
    """
    with open(script_path, 'r', encoding='utf-8') as f:
        slides = json.load(f)

    gen = PptxGenerator(mode=mode, bg_dir=bg_dir)
    slide_objects = []  # track (slide_obj, audio_path)

    for s in slides:
        t = s.get("type", "content")
        note = s.get("note", "")
        audio = s.get("audio", "")
        if audio and audio_dir and not os.path.isabs(audio):
            audio = os.path.join(audio_dir, audio)
        img = s.get("image", "")

        if t == "title":
            sl = gen.add_title_slide(s.get("title",""), s.get("subtitle",""), note)
        elif t == "objectives":
            sl = gen.add_objectives_slide(s.get("title",""), s.get("items",[]), note)
        elif t == "agenda":
            sl = gen.add_agenda_slide(s.get("title",""), s.get("items",[]), note)
        elif t == "divider":
            sl = gen.add_divider_slide(s.get("section_num",1), s.get("section_title",""), note)
        elif t == "content":
            sl = gen.add_content_slide(s.get("title",""), s.get("bullets",[]), img or None, note)
        elif t == "summary":
            sl = gen.add_summary_slide(s.get("title",""), s.get("items",[]), note)
        elif t == "exercise":
            sl = gen.add_exercise_slide(s.get("title",""), s.get("items",[]), note)
        elif t == "closing":
            sl = gen.add_closing_slide(s.get("title",""), s.get("next_lesson",""), note)
        else:
            sl = gen.add_content_slide(s.get("title",""), s.get("bullets",[]), None, note)

        slide_objects.append((sl, audio))

    # Inject audio
    for sl, audio in slide_objects:
        if audio and os.path.exists(audio):
            gen.add_audio(sl, audio)

    gen.save(output_path)
    return output_path


# ============================================================
# QUICK TEST
# ============================================================
if __name__ == "__main__":
    bg_dir = os.path.join(os.path.dirname(__file__), '..', '..', '..', '..', 'pptx', 'temp_background')
    bg_dir = os.path.normpath(bg_dir)
    out_dir = os.path.normpath(os.path.join(os.path.dirname(__file__), '..', '..', '..', '..', 'pptx'))

    for mode in ["bg1", "bg2"]:
        gen = PptxGenerator(mode=mode, bg_dir=bg_dir)
        gen.add_title_slide("Cong cu va Moi truong", "Phan I - Bai 02 - Nhap Mon Hoc May")
        gen.add_objectives_slide(["Thiet lap Google Colab", "Thao tac NumPy", "Doc du lieu Pandas", "Truc quan hoa Matplotlib"])
        gen.add_agenda_slide(["Moi truong tinh toan", "He sinh thai thu vien"])
        gen.add_divider_slide(1, "Moi truong Tinh toan Tuong tac")
        gen.add_content_slide("Google Colab", ["Jupyter Notebook tren cloud", "Mien phi GPU", "Chia se de dang"])
        gen.add_content_slide("Quan ly thu vien", [
            {"point": "import numpy as np", "description": "Thu vien tinh toan so hoc"},
            {"point": "import pandas as pd", "description": "Thu vien xu ly du lieu bang"},
        ])
        gen.add_divider_slide(2, "He sinh thai Thu vien Nen tang")
        gen.add_content_slide("NumPy", ["Mang ndarray dong nhat", "Vector hoa nhanh gap 100x", "Broadcasting"])
        gen.add_summary_slide(["Google Colab = moi truong sandbox", "NumPy = tinh toan nhanh", "Pandas = du lieu bang", "Matplotlib = truc quan hoa"])
        gen.add_exercise_slide(["Tao vector NumPy 10 phan tu", "Nap California Housing", "Bao cao EDA"])
        gen.add_closing_slide("Tien xu ly Du lieu & EDA")
        out = os.path.join(out_dir, f"test_{mode}.pptx")
        gen.save(out)

    print("\nDone! Check pptx/test_bg1.pptx and pptx/test_bg2.pptx")
